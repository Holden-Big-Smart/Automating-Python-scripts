import os
import pandas as pd
from datetime import datetime
from pptx import Presentation
from difflib import SequenceMatcher

# === 文件路径设置 ===
csv_path = "屯門婦聯 - 會員及課程管理系統 - 課程.csv"
ppt_folder = "海报"
output_log = []

# === 关键列标题 ===
KEY_COLUMNS = {
    "名稱": None,
    "上課日期": None,
    "時間": None,
    "逢星期": None
}

# === 函数：提取 PPT 中指定文本框 ===
def replace_text_in_shape(prs, placeholder_name, new_text):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == placeholder_name:
                text_frame = shape.text_frame
                # 如果已有段落，则只修改第一个段落中的第一个run的内容，保留其格式
                if text_frame.paragraphs:
                    p = text_frame.paragraphs[0]
                    if p.runs:
                        run = p.runs[0]
                        run.text = new_text
                    else:
                        # 没有run就创建一个新的 run（会丢失格式）
                        p.text = new_text
                else:
                    # 如果没有段落就设置整个 text（会丢失格式）
                    text_frame.text = new_text
                return True
    return False


# === 函数：处理日期格式 ===
def format_date(raw_date):
    try:
        start_end = raw_date.split("|")
        start = datetime.strptime(start_end[0].split(" ")[0], "%Y-%m-%d")
        end = datetime.strptime(start_end[1].split(" ")[0], "%Y-%m-%d")
        return f"日期 : {start.day:02d}/{start.month:02d}-{end.day:02d}/{end.month:02d}/{end.year}"
    except Exception as e:
        return "日期格式错误"

# === 函数：处理时间格式 ===
def format_time(raw_time, raw_weekday):
    try:
        start_end = raw_time.split("|")
        start = start_end[0].split(" ")[0]
        end = start_end[1].split(" ")[0]
        weekday = f"({raw_weekday.strip()})" if raw_weekday else ""
        return f"時間 : {start}-{end}{weekday}"
    except Exception as e:
        return "時間格式错误"

# === 函数：模糊匹配课程名称 ===
def get_best_match(course_name, name_list):
    best_score = 0
    best_index = -1
    for i, name in enumerate(name_list):
        score = SequenceMatcher(None, course_name, str(name)).ratio()
        if score > best_score:
            best_score = score
            best_index = i
    return best_index, best_score

# === 加载 CSV 数据并检查关键列 ===
try:
    df = pd.read_csv(csv_path)
except Exception as e:
    print(f"❌ 无法读取课程数据文件：{csv_path}")
    input("按 Enter 退出...")
    exit()

# 确保关键列存在
missing_columns = [col for col in KEY_COLUMNS if col not in df.columns]
if missing_columns:
    print("❌ 缺少以下关键列：", ", ".join(missing_columns))
    input("按 Enter 退出...")
    exit()

# 提取列索引（防止列顺序偏移）
for col in KEY_COLUMNS:
    KEY_COLUMNS[col] = df.columns.get_loc(col)

# === 遍历所有 PPT 文件 ===
for filename in os.listdir(ppt_folder):
    if not filename.endswith(".pptx"):
        continue

    ppt_path = os.path.join(ppt_folder, filename)
    course_name = os.path.splitext(filename)[0]

    match_index, match_score = get_best_match(course_name, df["名稱"])
    if match_index == -1:
        log = f"❗ 未找到与「{course_name}」匹配的课程"
        output_log.append(log)
        print(log)
        continue

    # 获取匹配行数据
    row = df.iloc[match_index]
    raw_date = row["上課日期"]
    raw_time = row["時間"]
    raw_weekday = row["逢星期"]

    # 格式转换
    formatted_date = format_date(raw_date)
    formatted_time = format_time(raw_time, raw_weekday)

    # 打开 PPT 并替换
    try:
        prs = Presentation(ppt_path)
        date_result = replace_text_in_shape(prs, "日期", formatted_date)
        time_result = replace_text_in_shape(prs, "時間", formatted_time)

        if not (date_result and time_result):
            log = f"⚠️ 未在「{filename}」中找到“日期”或“時間”文本框"
        else:
            prs.save(ppt_path)
            log = f"✅ 已更新「{filename}」：{formatted_date} / {formatted_time}"

        output_log.append(log)
        print(log)
    except Exception as e:
        log = f"❌ 处理「{filename}」时出错：{e}"
        output_log.append(log)
        print(log)

# === 完成提示 ===
print("\n--- 脚本执行完毕 ---")
input("按 Enter 退出...")
