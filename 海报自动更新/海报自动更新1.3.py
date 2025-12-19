# 先追加再删除

import os
import pandas as pd
from datetime import datetime
from pptx import Presentation
from difflib import SequenceMatcher

# === 文件路径设置 ===
csv_path = "屯門婦聯 - 會員及課程管理系統 - 課程.csv"
ppt_folder = "海报"
log_path = os.path.join(ppt_folder, "更新日誌.txt")

# === 关键列标题 ===
KEY_COLUMNS = {
    "名稱": None,
    "上課日期": None,
    "時間": None,
    "逢星期": None
}

# === 先追加后删除 替换文本，保留样式 ===
def replace_text_in_shape(prs, placeholder_name, new_text):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == placeholder_name:
                text_frame = shape.text_frame
                if not text_frame.paragraphs:
                    return False

                para = text_frame.paragraphs[0]
                if not para.runs:
                    para.text = new_text
                    return True

                # ✅ 第一步：保留第一个 run，删除其余 run（用 _p.remove）
                first_run = para.runs[0]
                for run in list(para.runs[1:]):
                    para._p.remove(run._r)

                # ✅ 第二步：追加新文本内容到第一个 run 后面
                first_run.text += new_text

                # ✅ 第三步：删除旧内容（保留的 run 的开头部分）
                first_run.text = first_run.text[-len(new_text):]

                return True
    return False


# === 函数：处理日期格式 ===
def format_date(raw_date):
    try:
        start_end = raw_date.split("|")
        start = datetime.strptime(start_end[0].split(" ")[0], "%Y-%m-%d")
        end = datetime.strptime(start_end[1].split(" ")[0], "%Y-%m-%d")
        return f"日期 : {start.day:02d}/{start.month:02d}-{end.day:02d}/{end.month:02d}/{end.year}"
    except:
        return "日期格式錯誤"

# === 函数：处理时间格式 ===
def format_time(raw_time, raw_weekday):
    try:
        start_end = raw_time.split("|")
        start = start_end[0].split(" ")[0]
        end = start_end[1].split(" ")[0]
        weekday = f"({raw_weekday.strip()})" if raw_weekday else ""
        return f"時間 : {start}-{end}{weekday}"
    except:
        return "時間格式錯誤"

# === 匹配函数：模糊匹配文件名到课程名称列 ===
def get_best_match(course_name, name_list):
    course_name_lower = course_name.lower()

    best_score = 0
    best_index = -1

    for i, raw_name in enumerate(name_list):
        name = str(raw_name)
        name_lower = name.lower()

        if name_lower == course_name_lower:
            return i, 2.0  # 完全匹配，直接返回

        score = SequenceMatcher(None, course_name_lower, name_lower).ratio()
        # PPT名称包含在课程名称中：
        if course_name_lower in name_lower:
            score += 0.2
        # 课程名称以PPT名称开头：
        if name_lower.startswith(course_name_lower):
            score += 0.2

        if score > best_score:
            best_score = score
            best_index = i

    return best_index, best_score

# === 加载 CSV 数据并检查关键列 ===
try:
    df = pd.read_csv(csv_path)
except Exception as e:
    print(f"❌ 無法讀取課程資料文件：{csv_path}")
    input("按 Enter 鍵退出...")
    exit()

missing_columns = [col for col in KEY_COLUMNS if col not in df.columns]
if missing_columns:
    print("❌ 缺少以下關鍵列：", ", ".join(missing_columns))
    input("按 Enter 鍵退出...")
    exit()

for col in KEY_COLUMNS:
    KEY_COLUMNS[col] = df.columns.get_loc(col)


with open(log_path, "w", encoding="utf-8") as f_log:
    f_log.write("【海報自動更新日誌】\n\n")

# === 遍历 PPTX 文件 ===
for filename in os.listdir(ppt_folder):
    if not filename.endswith(".pptx"):
        continue

    ppt_path = os.path.join(ppt_folder, filename)
    course_name = os.path.splitext(filename)[0]

    match_index, match_score = get_best_match(course_name, df["名稱"])
    if match_index == -1:
        log = f"❗ 無法匹配「{course_name}」"
        print(log)
        with open(log_path, "a", encoding="utf-8") as f_log:
            f_log.write(log + "\n")
        continue

    row = df.iloc[match_index]
    raw_title = str(row["名稱"])
    raw_date = row["上課日期"]
    raw_time = row["時間"]
    raw_weekday = row["逢星期"]

    formatted_date = format_date(raw_date)
    formatted_time = format_time(raw_time, raw_weekday)

    try:
        prs = Presentation(ppt_path)
        date_success = replace_text_in_shape(prs, "日期", formatted_date)
        time_success = replace_text_in_shape(prs, "時間", formatted_time)

        if not (date_success and time_success):
            log = f"⚠️ 「{filename}」未找到「日期」或「時間」文本框"
        else:
            prs.save(ppt_path)
            log = f"「{filename}」 ⇒ 「{raw_title}」\n{formatted_date} / {formatted_time}"

        print(log)
        with open(log_path, "a", encoding="utf-8") as f_log:
            f_log.write(log + "\n\n")
    except Exception as e:
        log = f"❌ 處理「{filename}」時出錯：{e}"
        print(log)
        with open(log_path, "a", encoding="utf-8") as f_log:
            f_log.write(log + "\n\n")

print("\n✅ 全部處理完成，日誌已儲存在：", log_path)
input("請按 Enter 鍵退出...")
