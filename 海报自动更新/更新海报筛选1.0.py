import os
import shutil
import re
from datetime import datetime
from pptx import Presentation

# ===== 路径设置 =====
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DISABLED_DIR = os.path.join(BASE_DIR, "海报", "Disabled")
OUTPUT_DIR = os.path.join(BASE_DIR, "海报")
LOG_PATH = os.path.join(OUTPUT_DIR, "筛选日志.txt")


def parse_date_text(date_text: str):
    """
    从字符串中提取开始日期（增强版，兼容多种格式）
    支持：
        日期：04/12-05/02/2026
        日期: 04/12-05/02/2026
        日期 : 04/12-05/02/2026
        日期:04/12-05/02/2026
    """
    pattern = r"日期\s*[:：]\s*(\d{2}/\d{2})-(\d{2}/\d{2}/\d{4})"
    match = re.search(pattern, date_text)

    if not match:
        raise ValueError(f"无法解析日期字段：{date_text}")

    start_str = match.group(1)  # 如 18/11
    end_str = match.group(2)    # 如 09/12/2025

    start_day, start_month = map(int, start_str.split("/"))
    end_day, end_month, end_year = map(int, end_str.split("/"))

    # 判断跨年
    if start_month > end_month:
        start_year = end_year - 1
    else:
        start_year = end_year

    start_date = datetime(start_year, start_month, start_day)

    return start_date, f"{start_str}-{end_str}"


def process_pptx_file(filepath, today):
    filename = os.path.basename(filepath)
    prs = Presentation(filepath)

    date_text = None

    # ===== 查找名为“日期”的 shape =====
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name.strip() == "日期" and shape.has_text_frame:
                text = shape.text.strip()

                # 放宽判断，只要包含关键词“日期”即可
                if "日期" in text:
                    date_text = text
                    break

        if date_text:
            break

    if not date_text:
        print(f"⚠️ 未找到日期字段：{filename}")
        return

    try:
        start_date, raw_date_str = parse_date_text(date_text)
        diff_days = (today - start_date).days

        if diff_days >= 0:
            # ===== 条件符合，复制文件 =====
            shutil.copy2(filepath, os.path.join(OUTPUT_DIR, filename))

            # ===== 写入日志 =====
            with open(LOG_PATH, "a", encoding="utf-8") as f:
                f.write(f"{filename}\n")
                f.write(f"日期：{raw_date_str}\n")
                f.write(f"课程已开始{diff_days}日(当前日期{today.strftime('%d/%m/%Y')})\n\n")

            print(f"✅ 已复制：{filename}（开始于 {start_date.strftime('%d/%m/%Y')}，已开始 {diff_days} 日）")

        else:
            print(f"⏭ 跳过未开始课程：{filename}")

    except Exception as e:
        print(f"❌ 处理文件时出错：{filename}，错误：{e}")


def main():
    print("📂 正在筛选海报...")
    today = datetime.today()

    pptx_files = [f for f in os.listdir(DISABLED_DIR) if f.lower().endswith(".pptx")]

    if not pptx_files:
        print("⚠️ 未找到任何 PPTX 文件")
        return

    for file in pptx_files:
        full_path = os.path.join(DISABLED_DIR, file)
        process_pptx_file(full_path, today)

    print("✅ 海报筛选完毕！")


if __name__ == "__main__":
    main()
